import streamlit as st
from streamlit_folium import st_folium
import folium
from geopy.geocoders import Nominatim
import pandas as pd
import numpy as np
import plotly.express as px
import matplotlib.pyplot as plt
from shapely.geometry import Point
import geopandas as gpd
import contextily as ctx
import pandas_gbq
from google.oauth2 import service_account

FR = "Français"
EN = "English"

if "local" not in st.session_state:
    st.session_state.local = FR  # Default local

# Sidebar Local Selection
with st.sidebar.expander("**Language**"):
    local = st.radio(
        "Select language / Sélectionnez la langue:",
        (FR, EN),
        index=0 if st.session_state.local == FR else 1,
    )

# Update session state when a new local is selected
if local != st.session_state.local:
    st.session_state.local = local


st.write("# Real Estate Snapshot - Reverse ETL")

st.write("## Introduction")

if local == EN:
    st.markdown(
        """
        Part of data analysis also involves transforming data and returning it to business teams in a clean and efficient way, so that they can accomplish their missions. It's what we call <b>Reverse ETL</b>.

        In this example, we developed an app for real estate agents. With this app, agents can choose a location in France and obtain real estate information within a radius they define.

        A button at the bottom of the page allows agents to download data in PowerPoint format, for offline access to information.
        """,
        unsafe_allow_html=True,
    )
else:
    st.markdown(
        """
        Une partie de la data analyse consiste également à transformer les données et à les restituer aux équipes métier d'une manière propre et efficace, afin qu'elles puissent accomplir leurs missions. C'est ce qu'on appelle le <b>Reverse ETL</b>.

        Dans cet exemple, nous avons développé une app pour des agents immobilier. Avec cette app, les agents peuvent choisir une localisation en France et obtenir les informations sur l'immobilier dans un rayon qu'ils définissent.

        Un bouton en bas de page permet aux agents de télécharger les données en PowerPoint, afin d'accéder aux informations hors ligne.
        """,
        unsafe_allow_html=True,
    )

st.divider()
dict_translation = {
    "map_search": {
        FR: "Saisissez une adresse ou un lieu :",
        EN: "Enter an address or location:",
    },
    "search": {FR: "Rechercher", EN: "Search"},
    "year_scorecard": {FR: "Sélectionnez une année", EN: "Select a year"},
    "distance": {
        FR: "Choisissez une distance (en mètres)",
        EN: "Choose distance (meters)",
    },
    "generate_snapshot": {
        FR: "Générer le snapshot",
        EN: "Generate Snapshot)",
    },
    "no_transaction": {
        FR: "Aucune vente n'a été trouvée dans la zone sélectionnée. Essayez d'augmenter le rayon ou de sélectionner un autre endroit.",
        EN: "No transactions found in the selected area. Try increasing the radius or selecting another location.",
    },
    "total_sales": {
        FR: "Nombre de ventes",
        EN: "Number of sales",
    },
    "avg_price": {
        FR: "Prix moyen de vente",
        EN: "Average sales price",
    },
    "avg_price_squaremetter": {
        FR: "Prix moyen du du m²",
        EN: "Average price of m²",
    },
    "avg_surface": {
        FR: "Surface moyenne",
        EN: "Average surface",
    },
    "house": {
        FR: "##### :house: Focus sur les maisons",
        EN: "##### :house: Focus on houses",
    },
    "flat": {
        FR: "##### :office: Focus sur les appartements",
        EN: "##### :office: Focus on flats",
    },
    "snapshot_year": {
        FR: "### :round_pushpin: Snapshot par année",
        EN: "### :round_pushpin: Snapshot per year",
    },
    "snapshot_general": {
        FR: "### :chart_with_upwards_trend: Vue d'ensemble",
        EN: "### :chart_with_upwards_trend: General overview",
    },
    "vefa": {
        FR: "Proportion de ventes sur plan (VEFA)",
        EN: "Proportion of off-plan sales (VEFA)",
    },
    "building_proportion": {
        FR: "Proportion par type de bâtiment",
        EN: "Distribution of Building Types",
    },
    "year": {
        FR: "Années",
        EN: "Years",
    },
    "building": {
        FR: "Type de bâtiment",
        EN: "Building type",
    },
    "filter_building": {
        FR: "Filtrer les graphiques suivants par type de bâtiment",
        EN: "Filter next charts by building type",
    },
    "house_filter": {
        FR: "Maison",
        EN: "House",
    },
    "flat_filter": {
        FR: "Appartement",
        EN: "Flat",
    },
    "dl_ppt": {
        FR: ":bar_chart: Télécharger le PowerPoint",
        EN: ":bar_chart: Download PowerPoint",
    },
    "mono_list_city": {
        FR: ":bar_chart: Télécharger le PowerPoint",
        EN: ":bar_chart: Download PowerPoint",
    },
}

revers_dict = {
    ":house: Maison": "Maison",
    ":office: Appartement": "Appartement",
    ":house: House": "Maison",
    ":office: Flat": "Appartement",
}


# Set default map center (France)
france_center = [46.603354, 1.888334]

# Initialize session state variables
if "last_location" not in st.session_state:
    st.session_state.last_location = None

if "radius" not in st.session_state:
    st.session_state.radius = 1000  # Default 1000m

if "snapshot_generated" not in st.session_state:
    st.session_state.snapshot_generated = False

if local == FR:
    st.write("### Rechercher une localisation")
    st.markdown(
        """
                Vous pouvez indiquer le nom d'une ville, ou une adresse dans la barre de recherche ci-dessous.<br>
                Si vous le souhaitez, vous pouvez zoomer et cliquer directement sur la carte pour placer le curseur.<br>
                Enfin, utilisez le slider pour définir le rayon la zone de recherche autour du curseur (en mètres)<br>
                """,
        unsafe_allow_html=True,
    )
else:
    st.write("### Search for a location")
    st.markdown(
        """
                You can enter a city name or address in the search bar below.<br>
                If you wish, you can zoom in and click directly on the map to place the cursor.<br>
                Finally, use the slider to define the radius of the search zone around the cursor (in meters).<br>
                """,
        unsafe_allow_html=True,
    )

# Address search box
address = st.text_input(dict_translation["map_search"][local])

if st.button(dict_translation["search"][local]):
    if address:
        geolocator = Nominatim(user_agent="streamlit-app")
        location = geolocator.geocode(address)

        if location:
            st.session_state.last_location = [location.latitude, location.longitude]
            st.success(f"Location found: {location.latitude}, {location.longitude}")
        else:
            st.error("Location not found. Try a different address.")

# Distance selection slider

# Temporary variable for the slider (to avoid instant updates)
radius = st.slider(
    dict_translation["distance"][local], 100, 10000, st.session_state.radius
)

# Create Folium map
m = folium.Map(location=france_center, zoom_start=6)

# Add a marker and a circle around the selected location
if st.session_state.last_location:
    lat, lon = st.session_state.last_location

    # Marker
    folium.Marker(
        location=[lat, lon],
        popup=f"Lat: {lat}, Lng: {lon}",
        icon=folium.Icon(color="red"),
    ).add_to(m)

    # Circle with light red color
    folium.Circle(
        location=[lat, lon],
        radius=radius,  # Use the selected radius
        color="red",
        fill=True,
        fill_color="lightcoral",  # Light red fill
        fill_opacity=0.3,
    ).add_to(m)

    # Center map on selected location
    m.location = [lat, lon]
    m.zoom_start = 14

# Display the map and capture user interaction
map_data = st_folium(m, width=700, height=500)

# Update session state with the last clicked location
if map_data and "last_clicked" in map_data and map_data["last_clicked"]:
    st.session_state.last_location = [
        map_data["last_clicked"]["lat"],
        map_data["last_clicked"]["lng"],
    ]

if local == FR:
    st.markdown(
        "Une fois que vous avez placé le curseur sur la carte, cliquez sur le bouton ci-dessous pour générer le snapshot"
    )
else:
    st.markdown(
        "Once you have placed the cursor on the map, click on the button below to generate the snapshot."
    )

# Button to trigger snapshot generation
if st.button(dict_translation["generate_snapshot"][local]):
    if st.session_state.last_location:
        lat, lon = st.session_state.last_location  # Selected location from map

        # Load city data
        def get_square_coordinates(lat, lon, radius_meters):
            earth_radius = 6378137.0

            delta_lat = (radius_meters / earth_radius) * (180 / np.pi)
            delta_lon = (radius_meters / (earth_radius * np.cos(np.radians(lat)))) * (
                180 / np.pi
            )

            max_lat = lat + delta_lat
            min_lat = lat - delta_lat
            max_lon = lon + delta_lon
            min_lon = lon - delta_lon

            return max_lat, min_lat, max_lon, min_lon

        # Get square boundary
        max_lat, min_lat, max_lon, min_lon = get_square_coordinates(lat, lon, radius)

        # SQL Query
        sql_query = f"""
        SELECT * FROM dbt_mbennis.mart_city_snapshot 
        WHERE latitude BETWEEN {min_lat} AND {max_lat}
        AND longitude BETWEEN {min_lon} AND {max_lon}
        ORDER BY date_year;
        """

        def get_data(
            sql: str, project_id: str, credentials: service_account.Credentials
        ) -> pd.DataFrame:
            df = pandas_gbq.read_gbq(
                sql, project_id=project_id, credentials=credentials
            )
            return df

        #### Load credentials from Streamlit secrets
        credentials_dict = dict(st.secrets["bigquery"])
        credentials = service_account.Credentials.from_service_account_info(
            credentials_dict
        )
        project_id = "mimetic-surfer-402208"

        df_transac_snapshot = get_data(sql_query, project_id, credentials)

        df_transac_snapshot["date_year"] = df_transac_snapshot["date_year"].astype(int)

        # TEST DATA WHILE DESIGNING, TO BE REPLACE BY SQL QUERY
        # df_transac_snapshot_test = pd.read_csv(
        #     "/home/mehdibennis/projects/cv_streamlit/files/data/PrixImmobilier/plotly/dbt_analyses_explo.csv"
        # )

        # df_transac_snapshot = df_transac_snapshot[
        #     (df_transac_snapshot["latitude"].between(min_lat, max_lat))
        #     & (df_transac_snapshot["longitude"].between(min_lon, max_lon))
        # ]

        type_batiment_counts = df_transac_snapshot["type_batiment"].value_counts()
        vefa_counts = df_transac_snapshot["vefa"].value_counts()

        tcd_transac_total = pd.pivot_table(
            df_transac_snapshot,
            index="date_year",
            values=["id_transaction", "prix", "avg_price_squaredmetters"],
            aggfunc={
                "id_transaction": "count",
                "prix": "mean",
                "avg_price_squaredmetters": "mean",
            },
            fill_value=0,
        )

        tcd_transac_type_batiment = pd.pivot_table(
            df_transac_snapshot,
            index="date_year",
            columns="type_batiment",
            values=[
                "id_transaction",
                "prix",
                "surface_habitable",
                "avg_price_squaredmetters",
            ],
            aggfunc={
                "id_transaction": "count",
                "prix": "mean",
                "surface_habitable": "mean",
                "avg_price_squaredmetters": "mean",
            },
            fill_value=0,
        )

        # Ensure 'date_year' is sorted and unique
        all_years = df_transac_snapshot["date_year"].unique()
        all_types = df_transac_snapshot["type_batiment"].unique()

        # Create a complete MultiIndex
        multi_index = pd.MultiIndex.from_product(
            [all_years, all_types], names=["date_year", "type_batiment"]
        )

        df_grouped_plotly = (
            df_transac_snapshot.groupby(["date_year", "type_batiment"])
            .agg(
                {
                    "id_transaction": "count",
                    "prix": "mean",
                    "avg_price_squaredmetters": "mean",
                }
            )
            .reset_index()
        )

        df_grouped_plotly = (
            df_grouped_plotly.set_index(["date_year", "type_batiment"])
            .reindex(multi_index)
            .fillna(0)
            .reset_index()
        )

        df_grouped_plotly["id_transaction"] = df_grouped_plotly[
            "id_transaction"
        ].astype(int)
        df_grouped_plotly["prix"] = df_grouped_plotly["prix"].astype(int)
        df_grouped_plotly["avg_price_squaredmetters"] = df_grouped_plotly[
            "avg_price_squaredmetters"
        ].astype(int)

        list_cities = list(df_transac_snapshot["formated_cleaned_ville"].unique())
        formatted_list_cities = ", ".join(map(str, list_cities))

        if df_transac_snapshot.empty:
            st.error(dict_translation["no_transaction"][local])
            st.stop()  # Stops execution to prevent Key

        # Store results in session state
        st.session_state.snapshot_generated = True
        st.session_state.formatted_list_cities = formatted_list_cities
        st.session_state.list_cities = list_cities
        st.session_state.type_batiment_counts = type_batiment_counts
        st.session_state.vefa_counts = vefa_counts
        st.session_state.tcd_transac_total = tcd_transac_total
        st.session_state.tcd_transac_type_batiment = tcd_transac_type_batiment
        st.session_state.df_grouped_plotly = df_grouped_plotly


# Display snapshot UI if data is generated
if st.session_state.snapshot_generated:
    # Retrieve data from session state
    formatted_list_cities = st.session_state.formatted_list_cities
    type_batiment_counts = st.session_state.type_batiment_counts
    vefa_counts = st.session_state.vefa_counts
    tcd_transac_total = st.session_state.tcd_transac_total
    tcd_transac_type_batiment = st.session_state.tcd_transac_type_batiment
    df_grouped_plotly = st.session_state.df_grouped_plotly
    list_cities = st.session_state.list_cities

    st.write(dict_translation["snapshot_year"][local])

    year = st.selectbox(
        dict_translation["year_scorecard"][local],
        (2024, 2023, 2022, 2021, 2020, 2019, 2018, 2017, 2016, 2015, 2014),
    )

    def safe_get_value_total(df, value_col, year):
        """Retrieve value safely, return 0 if column or value is missing for tcd transac total"""
        if value_col in df.columns:
            return int(np.nan_to_num(df.get(value_col, {}).get(year, 0), nan=0))
        return 0

    def safe_get_value_houseflat(df, value_col, type_col, year):
        """Retrieve value safely, return 0 if column or value is missing for tcd transac type batiment"""
        if type_col in df.columns.levels[1]:
            return int(
                np.nan_to_num(df[value_col].get(type_col, {}).get(year, 0), nan=0)
            )
        return 0

    # nb_sales = safe_get_value_total(tcd_transac_total, "id_transaction", year)
    nb_sales = safe_get_value_total(tcd_transac_total, "id_transaction", year)
    avg_price = safe_get_value_total(tcd_transac_total, "prix", year)
    avg_price_squaredmetter = safe_get_value_total(
        tcd_transac_total, "avg_price_squaredmetters", year
    )

    house_nb_sales = safe_get_value_houseflat(
        tcd_transac_type_batiment, "id_transaction", "Maison", year
    )
    house_avg_price = safe_get_value_houseflat(
        tcd_transac_type_batiment, "prix", "Maison", year
    )
    house_avg_price_squaredmetter = safe_get_value_houseflat(
        tcd_transac_type_batiment, "avg_price_squaredmetters", "Maison", year
    )
    house_avg_surface_habitable = safe_get_value_houseflat(
        tcd_transac_type_batiment, "surface_habitable", "Maison", year
    )

    flat_nb_sales = safe_get_value_houseflat(
        tcd_transac_type_batiment, "id_transaction", "Appartement", year
    )
    flat_avg_price = safe_get_value_houseflat(
        tcd_transac_type_batiment, "prix", "Appartement", year
    )
    flat_avg_price_squaredmetter = safe_get_value_houseflat(
        tcd_transac_type_batiment,
        "avg_price_squaredmetters",
        "Appartement",
        year,
    )
    flat_avg_surface_habitable = safe_get_value_houseflat(
        tcd_transac_type_batiment, "surface_habitable", "Appartement", year
    )

    dico_list_ville = {
        "mono": {
            FR: f"""
            Les ventes ont eu lieu dans la ville suivante :<b> {formatted_list_cities}</b> en <b>{year}</b>.
            """,
            EN: f"""
            Sales were located in this city : <b>{formatted_list_cities}</b> in <b>{year}</b>.
            """,
        },
        "multi": {
            FR: f"""
            Les ventes ont eu lieu dans les villes suivantes :<b> {formatted_list_cities}</b> en <b>{year}</b>.
            """,
            EN: f"""
            Sales were located in theses cities : <b>{formatted_list_cities}</b> in <b>{year}</b>.
            """,
        },
    }

    if len(list_cities) > 1:
        st.write(dico_list_ville["multi"][local], unsafe_allow_html=True)
    else:
        st.write(dico_list_ville["mono"][local], unsafe_allow_html=True)

    st.markdown(
        """
            <style>
            [data-testid="stMetricValue"] {
                font-size: 32px;
            }
            </style>
            """,
        unsafe_allow_html=True,
    )

    st.write("##### Total")
    col1, col2, col3 = st.columns(3)
    col1.metric(dict_translation["total_sales"][local], nb_sales, border=True)
    col2.metric(dict_translation["avg_price"][local], f"{avg_price}€", border=True)
    col3.metric(
        dict_translation["avg_price_squaremetter"][local],
        f"{avg_price_squaredmetter}€",
        border=True,
    )

    st.write(dict_translation["house"][local])

    col1, col2, col3, col4 = st.columns(4)
    col1.metric(dict_translation["total_sales"][local], house_nb_sales, border=True)
    col2.metric(
        dict_translation["avg_price"][local], f"{house_avg_price}€", border=True
    )
    col3.metric(
        dict_translation["avg_price_squaremetter"][local],
        f"{house_avg_price_squaredmetter}€",
        border=True,
    )
    col4.metric(
        dict_translation["avg_surface"][local],
        f"{house_avg_surface_habitable}m²",
        border=True,
    )

    st.write(dict_translation["flat"][local])

    col1, col2, col3, col4 = st.columns(4)
    col1.metric(dict_translation["total_sales"][local], flat_nb_sales, border=True)
    col2.metric(dict_translation["avg_price"][local], f"{flat_avg_price}€", border=True)
    col3.metric(
        dict_translation["avg_price_squaremetter"][local],
        f"{flat_avg_price_squaredmetter}€",
        border=True,
    )
    col4.metric(
        dict_translation["avg_surface"][local],
        f"{flat_avg_surface_habitable}m²",
        border=True,
    )
    st.divider()
    st.write(dict_translation["snapshot_general"][local])

    col1, col2 = st.columns(2, vertical_alignment="top")
    with col1:
        fig_type_batiment = px.pie(
            names=type_batiment_counts.index,  # Building types
            values=type_batiment_counts.values,  # Count of each type
            title=dict_translation["building_proportion"][local],
        )

        st.plotly_chart(fig_type_batiment)

    with col2:
        fig_vefa = px.pie(
            names=["Non VEFA", "VEFA"],
            values=[vefa_counts.get(0, 0), vefa_counts.get(1, 0)],
            title=dict_translation["vefa"][local],
        )
        st.plotly_chart(fig_vefa)

    radio = st.radio(
        dict_translation["filter_building"][local],
        (
            ":house: " + dict_translation["house_filter"][local],
            ":office: " + dict_translation["flat_filter"][local],
        ),
    )

    if radio == ":house: " + dict_translation["house_filter"][local]:
        chart_title = dict_translation["house_filter"][local]
    else:
        chart_title = dict_translation["flat_filter"][local]

    df_fitlered_plotly = df_grouped_plotly[
        df_grouped_plotly["type_batiment"] == revers_dict[radio]
    ]

    fig2 = px.line(
        df_fitlered_plotly, x="date_year", y="id_transaction", text="id_transaction"
    )

    fig2.update_layout(
        yaxis_title=dict_translation["total_sales"][local],
        xaxis_title=dict_translation["year"][local],
        title=chart_title + " : " + dict_translation["total_sales"][local],
    )
    fig2.update_traces(line=dict(color="red"), textposition="top center")

    st.plotly_chart(fig2)

    fig3 = px.line(df_fitlered_plotly, x="date_year", y="prix", text="prix")
    fig3.update_traces(line=dict(color="green"), textposition="top center")

    fig3.update_layout(
        yaxis_title=dict_translation["avg_price"][local],
        xaxis_title=dict_translation["year"][local],
        title=chart_title + " : " + dict_translation["avg_price"][local],
    )

    st.plotly_chart(fig3)

    fig4 = px.line(
        df_fitlered_plotly,
        x="date_year",
        y="avg_price_squaredmetters",
        text="avg_price_squaredmetters",
    )
    fig4.update_traces(line=dict(color="blue"), textposition="top center")

    fig4.update_layout(
        yaxis_title=dict_translation["avg_price_squaremetter"][local],
        xaxis_title=dict_translation["year"][local],
        title=chart_title + " : " + dict_translation["avg_price_squaremetter"][local],
    )

    st.plotly_chart(fig4)

    if local == FR:
        st.write("### Téléchargement PowerPoint")

        st.markdown(
            "Cliquez sur le bouton ci-dessous pour télécharger le snapshot au format PowerPoint."
        )
    else:
        st.write("### PowerPoint Download")

        st.markdown(
            "Click on the button below to download the snapshot in PowerPoint format."
        )

    from pptx import Presentation
    from io import BytesIO
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import (
        XL_CHART_TYPE,
        XL_LEGEND_POSITION,
        XL_TICK_LABEL_POSITION,
    )

    # # Create a PowerPoint presentation

    prs = Presentation("./files/ppt_template/CitySnapshot_template.pptx")

    ### SLIDE 1 - MAP

    slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if local == FR:
        title.text = "Localisation du snapshot"
    else:
        title.text = "Snapshot localisation"
        # Create a point geometry

    # CREATE MAP
    point = Point(lon, lat)
    radius_degrees = radius / 111320
    circle = point.buffer(radius_degrees)
    gdf = gpd.GeoDataFrame(geometry=[circle], crs="EPSG:4326").to_crs(epsg=3857)
    fig_map, ax = plt.subplots(figsize=(8, 8))
    gdf.plot(ax=ax, edgecolor="red", linewidth=2, facecolor="none")
    ctx.add_basemap(ax, source=ctx.providers.OpenStreetMap.Mapnik)
    ax.set_xticks([])
    ax.set_yticks([])

    # Save figure to an in-memory buffer
    image_buffer = BytesIO()
    plt.savefig(image_buffer, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig_map)  # Close figure to free memory
    image_buffer.seek(0)  # Move cursor to the start

    map_pic = slide.shapes.add_picture(
        image_buffer, Inches(1.5), Inches(1.5), Inches(2.81), Inches(3.86)
    )

    # Define text content
    if local == FR:
        text_content = f"Voici la carte du snapshot. Elle couvre les villes suivantes : {formatted_list_cities}"
    else:
        text_content = f"Here's the map on the snapshot. It covers the following cities : {formatted_list_cities}"
    # Add a text box beside the map
    text_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(3), Inches(4))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True  # Ensures text wraps within the box width
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centers text vertically

    # Add paragraph and style it
    p = text_frame.paragraphs[0]
    p.text = text_content
    p.alignment = PP_ALIGN.LEFT  # Align text to the left

    # Set font properties
    run = p.runs[0]
    run.font.name = "Montserrat"  # Font family
    run.font.size = Pt(12)  # Set a reasonable default size
    run.font.color.rgb = RGBColor(255, 255, 255)  # White text color

    # Enable auto-shrinking within PowerPoint (instead of `fit_text()`)
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    ### SLIDE 2 - Scorecards

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"Scorecard - {year}"

    rectangles = [
        # Line 1
        (1.5, 1.5, 2.5, 1),
        (4.13, 1.5, 2.5, 1),
        (6.75, 1.5, 2.5, 1),
        # Line 2
        (1.5, 2.75, 1.75, 1),
        (3.5, 2.75, 1.75, 1),
        (5.5, 2.75, 1.75, 1),
        (7.5, 2.75, 1.75, 1),
        # Line 3
        (1.5, 4, 1.75, 1),
        (3.5, 4, 1.75, 1),
        (5.5, 4, 1.75, 1),
        (7.5, 4, 1.75, 1),
    ]

    dico_scorecard = {
        FR: [
            {
                "metric": dict_translation["total_sales"][local] + " total",
                "score": nb_sales,
            },
            {
                "metric": dict_translation["avg_price"][local] + " total",
                "score": f"{avg_price}€",
            },
            {
                "metric": dict_translation["avg_price_squaremetter"][local] + " total",
                "score": f"{avg_price_squaredmetter}€",
            },
            ### Line 2
            {"metric": dict_translation["total_sales"][local], "score": house_nb_sales},
            {
                "metric": dict_translation["avg_price"][local],
                "score": f"{house_avg_price}€",
            },
            {
                "metric": dict_translation["avg_price_squaremetter"][local],
                "score": f"{house_avg_price_squaredmetter}€",
            },
            {
                "metric": dict_translation["avg_surface"][local],
                "score": f"{house_avg_surface_habitable}€",
            },
            ### Line 3
            {"metric": dict_translation["total_sales"][local], "score": flat_nb_sales},
            {
                "metric": dict_translation["avg_price"][local],
                "score": f"{flat_avg_price}€",
            },
            {
                "metric": dict_translation["avg_price_squaremetter"][local],
                "score": f"{flat_avg_price_squaredmetter}€",
            },
            {
                "metric": dict_translation["avg_surface"][local],
                "score": f"{flat_avg_surface_habitable}€",
            },
        ],
        EN: [
            {"metric": dict_translation["total_sales"][local], "score": nb_sales},
            {"metric": dict_translation["avg_price"][local], "score": f"{avg_price}€"},
            {
                "metric": dict_translation["avg_price_squaremetter"][local],
                "score": f"{avg_price_squaredmetter}€",
            },
            ### Line 2
            {"metric": dict_translation["total_sales"][local], "score": house_nb_sales},
            {
                "metric": dict_translation["avg_price"][local],
                "score": f"{house_avg_price}€",
            },
            {
                "metric": dict_translation["avg_price_squaremetter"][local],
                "score": f"{house_avg_price_squaredmetter}€",
            },
            {
                "metric": dict_translation["avg_surface"][local],
                "score": f"{house_avg_surface_habitable}€",
            },
            ### Line 3
            {"metric": dict_translation["total_sales"][local], "score": flat_nb_sales},
            {
                "metric": dict_translation["avg_price"][local],
                "score": f"{flat_avg_price}€",
            },
            {
                "metric": dict_translation["avg_price_squaremetter"][local],
                "score": f"{flat_avg_price_squaredmetter}€",
            },
            {
                "metric": dict_translation["avg_surface"][local],
                "score": f"{flat_avg_surface_habitable}€",
            },
        ],
    }

    rectangle_color = RGBColor(130, 199, 165)
    font_color = RGBColor(0, 0, 0)

    for i, (x, y, w, h) in enumerate(rectangles):
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(h)
        )  # Rectangle shape
        shape.fill.solid()
        shape.fill.fore_color.rgb = rectangle_color
        shape.line.fill.background()
        text_frame = shape.text_frame
        text_frame.clear()  # Clear any default text

        # First line of text
        p1 = text_frame.add_paragraph()
        p1.text = dico_scorecard[local][i]["metric"]
        p1.font.size = Pt(12)
        p1.font.color.rgb = font_color
        p1.alignment = PP_ALIGN.CENTER

        # Second line of text (bold, larger font)
        p2 = text_frame.add_paragraph()
        p2.text = str(dico_scorecard[local][i]["score"])
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = font_color
        p2.alignment = PP_ALIGN.CENTER

    house_pic = slide.shapes.add_picture(
        "./files/images/house.png", Inches(0.25), Inches(2.75), Inches(1), Inches(1)
    )

    flat_pic = slide.shapes.add_picture(
        "./files/images/flat.png", Inches(0.25), Inches(4), Inches(1), Inches(1)
    )

    ### SLIDE 3 - Proportions bulding types

    slide = prs.slides.add_slide(slide_layout)

    # Add title to the slide
    title = slide.shapes.title
    if local == FR:
        title.text = "Proportion par type de batiment & de vente sur plan (VEFA)"
    else:
        title.text = "Proportion of bulding type and off-plan sales (VEFA)"

    # Define chart position and size
    x, y, cx, cy = Inches(0.5), Inches(1.75), Inches(4), Inches(3.5)

    shape = slide.shapes.add_shape(1, x, y, cx, cy)  # Rectangle shape
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = ["Non VEFA", "VEFA"]
    chart_data.add_series("VEFA", (vefa_counts.get(0, 0), vefa_counts.get(1, 0)))

    # Add pie chart to the slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Define chart position and size
    x, y, cx, cy = Inches(5), Inches(1.75), Inches(4), Inches(3.5)

    shape = slide.shapes.add_shape(1, x, y, cx, cy)  # Rectangle shape
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = type_batiment_counts.index
    chart_data.add_series(
        dict_translation["building"][local], tuple(type_batiment_counts.values)
    )

    # Add pie chart to the slide
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    ################# Add line charts

    series_chart = [
        tcd_transac_total["id_transaction"],
        tcd_transac_total["prix"],
        tcd_transac_total["avg_price_squaredmetters"],
        tcd_transac_type_batiment["id_transaction"]["Maison"],
        tcd_transac_type_batiment["prix"]["Maison"],
        tcd_transac_type_batiment["avg_price_squaredmetters"]["Maison"],
        tcd_transac_type_batiment["id_transaction"]["Appartement"],
        tcd_transac_type_batiment["prix"]["Appartement"],
        tcd_transac_type_batiment["avg_price_squaredmetters"]["Appartement"],
    ]

    dict_chart = {
        FR: [
            "Nombre de ventes - Total",
            "Prix moyen de vente - Total",
            "Prix moyen du m² - Total",
            "Nombre de ventes - Maisons",
            "Prix moyen de vente - Maisons",
            "Prix moyen du m² - Maisons",
            "Nombre de ventes - Appartements",
            "Prix moyen de vente - Appartements",
            "Prix moyen du m² - Appartements",
        ],
        EN: [
            "Number of Sales - Total",
            "Average Price of Sales - Total",
            "Average Price of m² - Total",
            "Number of Sales - Houses",
            "Average Price of Sales - Houses",
            "Average Price of m² - Houses",
            "Number of Sales - Flats",
            "Average Price of Sales - Flats",
            "Average Price of m² - Flats",
        ],
    }

    for i in range(len(series_chart)):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = dict_chart[local][i]

        x, y, cx, cy = Inches(0.5), Inches(1.75), Inches(8.5), Inches(3.5)

        shape = slide.shapes.add_shape(1, x, y, cx, cy)  # Rectangle shape
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.fill.background()

        chart_data = CategoryChartData()
        chart_data.categories = series_chart[i].index
        chart_data.add_series(dict_chart[local][i], series_chart[i].to_list())

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            x,
            y,
            cx,
            cy,
            chart_data,
        ).chart

        chart.has_title = False
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.tick_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)

    # Save the presentation to a BytesIO buffer

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)  # Move the cursor to the beginning

    # Create a download button

    st.download_button(
        label=dict_translation["dl_ppt"][local],
        data=ppt_buffer,
        file_name="real_estate_stapshot.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
